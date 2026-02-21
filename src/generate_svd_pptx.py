from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # --- Config ---
    COLOR_PRIMARY = RGBColor(51, 51, 51)     # Dark Gray
    COLOR_ACCENT = RGBColor(212, 98, 42)    # Orange (Sunset)
    COLOR_GOLD = RGBColor(184, 134, 11)     # Gold
    COLOR_ALERT = RGBColor(196, 78, 90)     # Red
    COLOR_BLUE = RGBColor(70, 130, 180)     # Steel Blue
    COLOR_BG_TITLE = RGBColor(253, 243, 235) # Light Orange BG
    
    def set_font(run, size, bold=False, color=COLOR_PRIMARY):
        run.font.name = 'Arial' 
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_title_slide(prs, title_text, subtitle_text, badge_text):
        slide = prs.slides.add_slide(prs.slide_layouts[0]) 
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG_TITLE
        
        title = slide.shapes.title
        title.text = title_text
        set_font(title.text_frame.paragraphs[0].runs[0], 54, True, COLOR_ACCENT)
        
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        for run in subtitle.text_frame.paragraphs[0].runs:
            set_font(run, 24, False, COLOR_PRIMARY)
            
        # Badge
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(3)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        p = txBox.text_frame.paragraphs[0]
        p.text = badge_text
        set_font(p.runs[0], 14, True, COLOR_GOLD)

    def add_content_slide(prs, title_text, content_text):
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        title = slide.shapes.title
        title.text = title_text
        set_font(title.text_frame.paragraphs[0].runs[0], 36, True, COLOR_ACCENT)
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.word_wrap = True
        
        lines = content_text.strip().split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line: continue
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = line
            p.space_after = Pt(14)
            set_font(p.runs[0], 20, False, COLOR_PRIMARY)
        return slide

    def add_benchmark_chart_slide(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Chapter 05: ベンチマーク（実績見込）"
        set_font(title.text_frame.paragraphs[0].runs[0], 36, True, COLOR_ACCENT)

        # Draw Chart Area
        left_margin = Inches(1)
        top_start = Inches(2)
        bar_height = Inches(0.6)
        gap = Inches(0.4)
        max_width = Inches(6) # 100% width reference (roughly 30%)
        
        data = [
            ("ホテルF&B (Global)", 28.0, COLOR_PRIMARY, "28%"),
            ("独立系FSR (Japan)", 4.0, COLOR_BLUE, "4%"),
            ("SVD (R7実績見込)", 5.6, COLOR_ACCENT, "5.6%"),
            ("SVD FL比率", 73.0, COLOR_ALERT, "73.0% (危)")
        ]

        for i, (label, value, color, text_val) in enumerate(data):
            top = top_start + i * (bar_height + gap)
            
            # Label
            txBox = slide.shapes.add_textbox(left_margin, top, Inches(2.5), bar_height)
            p = txBox.text_frame.paragraphs[0]
            p.text = label
            set_font(p.runs[0], 16, True, COLOR_PRIMARY)
            
            # Bar (Scale: 100% = 7 inches? No, let's say 80% is max width)
            # Scaling: 80% = 6 inches -> 1% = 0.075 inches
            if "FL" in label:
                # FL is huge, scale differently or break visual
                width = Inches(6) # Max out
            else:
                width = Inches(value * 0.15) # 1% = 0.15 inch
                
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_margin + Inches(2.5), top, width, bar_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background() # No border

            # Value Label
            txBoxVal = slide.shapes.add_textbox(left_margin + Inches(2.5) + width + Inches(0.1), top, Inches(1.5), bar_height)
            p = txBoxVal.text_frame.paragraphs[0]
            p.text = text_val
            set_font(p.runs[0], 16, True, color)

        # Insight Text
        txBox = slide.shapes.add_textbox(left_margin, top_start + Inches(4.5), Inches(8), Inches(1))
        p = txBox.text_frame.paragraphs[0]
        p.text = "FL比率 73% の異常値が、DOP 5.6% の主因。\nここをコントロールすることこそが、最大のレバーである。"
        set_font(p.runs[0], 18, False, COLOR_ALERT)

    def add_usali_diagram_slide(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Chapter 03 & 04: USALI構造とDOP"
        set_font(title.text_frame.paragraphs[0].runs[0], 36, True, COLOR_ACCENT)

        cx = Inches(4.5) # Center X
        width = Inches(4)
        height = Inches(0.6)
        gap = Inches(0.2)
        start_y = Inches(1.8)

        # Layers
        layers = [
            ("売上高 (Revenue)", COLOR_PRIMARY),
            ("- 変動費 (Cost of Sales)", COLOR_PRIMARY),
            ("= 限界利益 (Marginal Profit)", COLOR_BLUE),
            ("- 個別固定費 (Labor, Direct Exp)", COLOR_PRIMARY),
            ("= DOP (貢献利益)", COLOR_ACCENT), # HERO
            ("- 共通固定費 (Undistributed)", COLOR_PRIMARY),
            ("= GOP (営業総利益)", COLOR_GOLD),
            ("- 所有コスト (Fixed Charges)", COLOR_PRIMARY),
            ("= 純利益 (Net Profit)", COLOR_PRIMARY)
        ]

        for i, (text, color) in enumerate(layers):
            top = start_y + i * (height + gap)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            if "DOP" in text:
                shape.line.color.rgb = COLOR_ACCENT
                shape.line.width = Pt(3)
            else:
                shape.line.fill.background()
            
            shape.text_frame.text = text
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    set_font(p.runs[0], 16, True if "DOP" in text or "GOP" in text else False, color)

        # Annotation
        txBox = slide.shapes.add_textbox(Inches(7.2), start_y + Inches(2.5), Inches(2.5), Inches(2))
        p = txBox.text_frame.paragraphs[0]
        p.text = "← ここがSVDの戦場\n(運営チームの責任範囲)"
        set_font(p.runs[0], 14, True, COLOR_ACCENT)

    def add_equation_flow_slide(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Chapter 07: SVDの方程式フロー"
        set_font(title.text_frame.paragraphs[0].runs[0], 36, True, COLOR_ACCENT)
        
        # Flowchart: Peaks -> Talent -> Value -> DOP
        nodes = [
            ("需要の波\n(Peaks)", Inches(1)),
            ("適材配置\n(Talent)", Inches(3.5)),
            ("価値最大化\n(Value)", Inches(6)),
            ("人時DOP\n(Profit)", Inches(8.5))
        ]
        
        y = Inches(3)
        w = Inches(1.8)
        h = Inches(1.2)
        
        for i, (label, x) in enumerate(nodes):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_ACCENT if i == 3 else RGBColor(240, 240, 240)
            shape.line.color.rgb = COLOR_ACCENT
            
            shape.text_frame.text = label
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    set_font(p.runs[0], 16, True, RGBColor(255,255,255) if i==3 else COLOR_PRIMARY)

            # Arrow
            if i < len(nodes) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Inches(0.1), y + h/2 - Inches(0.15), Inches(0.5), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = COLOR_PRIMARY
                arrow.line.fill.background()

        # Explanation
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
        p = txBox.text_frame.paragraphs[0]
        p.text = "MPで「波」を予測し、TIで「人」を配置する。\n無駄なアイドルタイムをゼロにすることで、\n投入工数あたりの価値（人時DOP）を最大化する。"
        set_font(p.runs[0], 20, False, COLOR_PRIMARY)

    # --- EXECUTE ---
    add_title_slide(prs, "DOPで測る。\nSVDの実力。", "札幌振興公社の料飲事業として\n運営力を可視化する独自指標", "SVD Intelligence")
    add_content_slide(prs, "Chapter 01: 結論から。", "SVDは札幌振興公社の料飲事業として観光地でレストランを展開するF&B事業者。\n\n■ 戦略的回答：\nSVDは「貢献利益（Contribution Margin）」で評価する。\n現場責任者がコントロール可能な収益と費用のみを評価対象とする「管理可能性原則」に基づく。")
    add_content_slide(prs, "Chapter 02: SVDの出店モデル", "■ 自社施設内出店（大倉山・藻岩山）\n家賃発生せず。固定費構造に強み。\n\n■ テナント出店（赤れんが・テレビ塔）\n商業施設等に賃借。変動費的コスト構造。\n\nいずれも「観光」が主戦場。")
    
    add_usali_diagram_slide(prs) # Visual Diagram 1
    
    add_content_slide(prs, "Chapter 04: SVD版 DOP の定義", "計算式：売上 − 原価 − 人件費 − 直接経費 = DOP\n※ 減価償却・家賃・金利は「所有コスト」として除外。\n\nDOPは「貢献利益」である。\n本社費や共通固定費を回収する源泉となる。")
    
    add_benchmark_chart_slide(prs) # Visual Chart 1
    
    add_content_slide(prs, "Chapter 05: 構造的差異の解説", "ホテルF&B（27%）とSVD（5.6%）の乖離理由は\n「コスト負担範囲（Cost Scope）」の違い。\n\nSVDは独立事業者として施設維持費も負担している。\n5.6%は極めて厳格な利益率だが、FL73%は改善余地がある。")
    add_content_slide(prs, "Chapter 06: DOP改善の5つのレバー", "1. Labor (TIによる人件費コントロール)\n2. Cost of Sales (原価率低減)\n3. Top-line (客単価・客数)\n4. Channel Mix (BG/TO/宴会)\n5. Energy (光熱費削減)")
    
    add_equation_flow_slide(prs) # Visual Flow 1
    
    add_content_slide(prs, "KPI: 人時DOP", "■ KPI: 人時DOP（DOP per Man-Hour）\n「1時間働いていくら利益を出したか」。\n売上だけでなく、利益ベースの生産性を追及する。\n\n■ Goal\n営業利益0.3%ではなく、DOP 7.3%（現状5.6%からの回復）。\nそしてこれを10%以上に引き上げること。")

    OUTPUT_FILE = '/Users/satoshiiga/Desktop/SVD_DOP_Presentation_Enhanced.pptx'
    prs.save(OUTPUT_FILE)
    print(f"Enhanced Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
